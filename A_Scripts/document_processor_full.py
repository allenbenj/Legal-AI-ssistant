import asyncio
import hashlib
import io
import logging
import mimetypes
import os
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Union
import json
import re

# Core dependencies
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# Structured data dependencies
try:
    import pandas as pd
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Markdown processing
try:
    import markdown
    from bs4 import BeautifulSoup
    MARKDOWN_AVAILABLE = True
except ImportError:
    MARKDOWN_AVAILABLE = False

from .base_agent import BaseAgent, AgentResult, AgentError, TaskPriority

logger = logging.getLogger(__name__)

class ProcessingStrategy:
    """Document processing strategy constants"""
    FULL_PROCESSING = "full_processing"      # Complete text extraction and analysis
    STRUCTURED_DATA = "structured_data"      # Extract to database tables
    REFERENCE_ONLY = "reference_only"        # Metadata + basic text only

class DocumentType:
    """Document type constants"""
    PDF = "pdf"
    DOCX = "docx"
    TXT = "txt"
    MD = "markdown"
    HTML = "html"
    RTF = "rtf"
    XLSX = "xlsx"
    XLS = "xls"
    CSV = "csv"
    PPTX = "pptx"
    PPT = "ppt"
    IMAGE = "image"
    UNKNOWN = "unknown"

class DocumentProcessorAgent(BaseAgent):
    """
    Enhanced agent for processing legal documents with intelligent strategy selection
    
    Processing Strategies:
    - FULL_PROCESSING: PDF, DOCX, TXT, MD, HTML, RTF - Complete analysis
    - STRUCTURED_DATA: XLSX, XLS, CSV - Extract to database tables
    - REFERENCE_ONLY: PPTX, PPT, Images - Metadata + basic reference
    """
    
    def __init__(self, services):
        super().__init__(services, "DocumentProcessor")
        
        # Check dependencies
        self._check_dependencies()
        
        # File type configuration with processing strategies
        self.file_config = {
            # Full processing files
            '.pdf': {'type': DocumentType.PDF, 'strategy': ProcessingStrategy.FULL_PROCESSING},
            '.docx': {'type': DocumentType.DOCX, 'strategy': ProcessingStrategy.FULL_PROCESSING},
            '.doc': {'type': DocumentType.DOCX, 'strategy': ProcessingStrategy.FULL_PROCESSING},
            '.txt': {'type': DocumentType.TXT, 'strategy': ProcessingStrategy.FULL_PROCESSING},
            '.md': {'type': DocumentType.MD, 'strategy': ProcessingStrategy.FULL_PROCESSING},
            '.html': {'type': DocumentType.HTML, 'strategy': ProcessingStrategy.FULL_PROCESSING},
            '.htm': {'type': DocumentType.HTML, 'strategy': ProcessingStrategy.FULL_PROCESSING},
            '.rtf': {'type': DocumentType.RTF, 'strategy': ProcessingStrategy.FULL_PROCESSING},
            
            # Structured data files
            '.xlsx': {'type': DocumentType.XLSX, 'strategy': ProcessingStrategy.STRUCTURED_DATA},
            '.xls': {'type': DocumentType.XLS, 'strategy': ProcessingStrategy.STRUCTURED_DATA},
            '.csv': {'type': DocumentType.CSV, 'strategy': ProcessingStrategy.STRUCTURED_DATA},
            
            # Reference only files
            '.pptx': {'type': DocumentType.PPTX, 'strategy': ProcessingStrategy.REFERENCE_ONLY},
            '.ppt': {'type': DocumentType.PPT, 'strategy': ProcessingStrategy.REFERENCE_ONLY},
            '.png': {'type': DocumentType.IMAGE, 'strategy': ProcessingStrategy.REFERENCE_ONLY},
            '.jpg': {'type': DocumentType.IMAGE, 'strategy': ProcessingStrategy.REFERENCE_ONLY},
            '.jpeg': {'type': DocumentType.IMAGE, 'strategy': ProcessingStrategy.REFERENCE_ONLY},
            '.tiff': {'type': DocumentType.IMAGE, 'strategy': ProcessingStrategy.REFERENCE_ONLY},
            '.bmp': {'type': DocumentType.IMAGE, 'strategy': ProcessingStrategy.REFERENCE_ONLY}
        }
        
        # Processing handlers by document type
        self.handlers = {
            DocumentType.PDF: self._process_pdf,
            DocumentType.DOCX: self._process_docx,
            DocumentType.TXT: self._process_txt,
            DocumentType.MD: self._process_markdown,
            DocumentType.HTML: self._process_html,
            DocumentType.RTF: self._process_rtf,
            DocumentType.XLSX: self._process_excel,
            DocumentType.XLS: self._process_excel,
            DocumentType.CSV: self._process_csv,
            DocumentType.PPTX: self._process_powerpoint,
            DocumentType.PPT: self._process_powerpoint,
            DocumentType.IMAGE: self._process_image_reference
        }
        
        # Legal document indicators
        self.legal_indicators = [
            'motion', 'complaint', 'affidavit', 'deposition', 'subpoena',
            'warrant', 'indictment', 'plea', 'sentence', 'judgment',
            'order', 'ruling', 'opinion', 'brief', 'memorandum',
            'constitution', 'amendment', 'statute', 'regulation',
            'violation', 'misconduct', 'evidence', 'testimony',
            'plaintiff', 'defendant', 'appellant', 'appellee'
        ]
    
    def _load_config(self) -> Dict[str, Any]:
        """Load document processor configuration"""
        return {
            'timeout': 600,  # 10 minutes for large documents
            'max_retries': 2,
            'retry_delay': 2.0,
            'max_file_size_mb': 100,
            'ocr_language': 'eng',
            'extract_images': True,
            'clean_text': True,
            'detect_duplicates': True,
            'preserve_structure': True,
            'extract_tables': True
        }
    
    def _check_dependencies(self) -> None:
        """Check if required dependencies are available"""
        deps_status = {
            'PyMuPDF (PDF)': PYMUPDF_AVAILABLE,
            'python-docx (DOCX)': DOCX_AVAILABLE,
            'OCR (Images)': OCR_AVAILABLE,
            'pandas/openpyxl (Excel)': EXCEL_AVAILABLE,
            'python-pptx (PowerPoint)': PPTX_AVAILABLE,
            'markdown/BeautifulSoup (MD/HTML)': MARKDOWN_AVAILABLE
        }
        
        for dep_name, available in deps_status.items():
            if not available:
                logger.warning(f"{dep_name} not available - related functionality disabled")
            else:
                logger.debug(f"{dep_name} available")
    
    async def _process_task(self, task_data: Any, metadata: Dict[str, Any]) -> Dict[str, Any]:
        """
        Process a document file with strategy-based handling
        
        Args:
            task_data: File path (str/Path) or file-like object
            metadata: Additional processing options
            
        Returns:
            Dictionary with extracted content and metadata based on processing strategy
        """
        # Validate input
        if isinstance(task_data, (str, Path)):
            file_path = Path(task_data)
            if not file_path.exists():
                raise AgentError(f"File not found: {file_path}", self.name)
        else:
            raise AgentError("Unsupported input type - expected file path", self.name)
        
        # Check file size
        file_size_mb = file_path.stat().st_size / (1024 * 1024)
        if file_size_mb > self.config['max_file_size_mb']:
            raise AgentError(f"File too large: {file_size_mb:.1f}MB (max: {self.config['max_file_size_mb']}MB)", self.name)
        
        # Determine document type and strategy
        file_info = self._detect_document_info(file_path)
        doc_type = file_info['type']
        strategy = file_info['strategy']
        
        if doc_type == DocumentType.UNKNOWN:
            raise AgentError(f"Unsupported file type: {file_path.suffix}", self.name)
        
        # Get handler
        handler = self.handlers.get(doc_type)
        if not handler:
            raise AgentError(f"No handler available for document type: {doc_type}", self.name)
        
        # Process the document
        logger.info(f"Processing {doc_type} document with {strategy} strategy: {file_path.name}")
        
        try:
            # Extract content based on strategy
            extraction_result = await handler(file_path, metadata)
            
            # Add common metadata
            result = {
                'file_path': str(file_path),
                'file_name': file_path.name,
                'file_size_bytes': file_path.stat().st_size,
                'document_type': doc_type,
                'processing_strategy': strategy,
                'processing_notes': extraction_result.get('notes', [])
            }
            
            # Strategy-specific result formatting
            if strategy == ProcessingStrategy.FULL_PROCESSING:
                result.update(self._format_full_processing_result(extraction_result, file_path))
            elif strategy == ProcessingStrategy.STRUCTURED_DATA:
                result.update(self._format_structured_data_result(extraction_result, file_path))
            elif strategy == ProcessingStrategy.REFERENCE_ONLY:
                result.update(self._format_reference_only_result(extraction_result, file_path))
            
            logger.info(f"Successfully processed {file_path.name} using {strategy} strategy")
            return result
            
        except Exception as e:
            logger.error(f"Error processing {file_path.name}: {e}")
            raise AgentError(f"Document processing failed: {e}", self.name)
    
    def _detect_document_info(self, file_path: Path) -> Dict[str, str]:
        """Detect document type and processing strategy"""
        extension = file_path.suffix.lower()
        
        if extension in self.file_config:
            return self.file_config[extension]
        
        # Try to detect from MIME type
        mime_type, _ = mimetypes.guess_type(str(file_path))
        if mime_type:
            # Fallback MIME type detection
            if mime_type.startswith('image/'):
                return {'type': DocumentType.IMAGE, 'strategy': ProcessingStrategy.REFERENCE_ONLY}
            elif mime_type == 'application/pdf':
                return {'type': DocumentType.PDF, 'strategy': ProcessingStrategy.FULL_PROCESSING}
            elif 'spreadsheet' in mime_type or 'excel' in mime_type:
                return {'type': DocumentType.XLSX, 'strategy': ProcessingStrategy.STRUCTURED_DATA}
            elif 'presentation' in mime_type or 'powerpoint' in mime_type:
                return {'type': DocumentType.PPTX, 'strategy': ProcessingStrategy.REFERENCE_ONLY}
            elif mime_type.startswith('text/'):
                return {'type': DocumentType.TXT, 'strategy': ProcessingStrategy.FULL_PROCESSING}
        
        return {'type': DocumentType.UNKNOWN, 'strategy': ProcessingStrategy.REFERENCE_ONLY}
    
    def _format_full_processing_result(self, extraction_result: Dict, file_path: Path) -> Dict[str, Any]:
        """Format result for full processing strategy"""
        text = extraction_result['text']
        
        # Generate document hash for duplicate detection
        content_hash = self._calculate_content_hash(text)
        
        # Detect if it's a legal document
        is_legal = self._is_legal_document(text)
        
        # Clean text if requested
        if self.config['clean_text']:
            text = self._clean_text(text)
        
        result = {
            'text': text,
            'content_hash': content_hash,
            'is_legal_document': is_legal,
            'metadata': extraction_result.get('metadata', {}),
            'page_count': extraction_result.get('page_count', 1),
            'word_count': len(text.split()),
            'character_count': len(text),
            'structure': extraction_result.get('structure', {})
        }
        
        # Add legal document classification if applicable
        if is_legal:
            result['legal_classification'] = self._classify_legal_document(text)
        
        return result
    
    def _format_structured_data_result(self, extraction_result: Dict, file_path: Path) -> Dict[str, Any]:
        """Format result for structured data strategy"""
        return {
            'tables': extraction_result.get('tables', []),
            'table_count': len(extraction_result.get('tables', [])),
            'total_rows': sum(len(table.get('data', [])) for table in extraction_result.get('tables', [])),
            'sheets': extraction_result.get('sheets', []),
            'metadata': extraction_result.get('metadata', {}),
            'database_schema': extraction_result.get('database_schema', {}),
            'text_content': extraction_result.get('text_content', ''),  # Any extracted text
            'is_legal_document': self._is_legal_document(extraction_result.get('text_content', ''))
        }
    
    def _format_reference_only_result(self, extraction_result: Dict, file_path: Path) -> Dict[str, Any]:
        """Format result for reference-only strategy"""
        return {
            'reference_type': 'file_reference',
            'quick_preview': extraction_result.get('preview_text', '')[:500],  # First 500 chars
            'metadata': extraction_result.get('metadata', {}),
            'slide_count': extraction_result.get('slide_count'),
            'image_count': extraction_result.get('image_count'),
            'accessible': True,  # File is accessible for future reference
            'summary': extraction_result.get('summary', '')
        }
    
    # === FULL PROCESSING HANDLERS ===
    
    async def _process_pdf(self, file_path: Path, metadata: Dict) -> Dict[str, Any]:
        """Process PDF document with full text extraction"""
        if not PYMUPDF_AVAILABLE:
            raise AgentError("PyMuPDF not available for PDF processing", self.name)
        
        try:
            doc = fitz.open(str(file_path))
            
            text_parts = []
            doc_metadata = {}
            notes = []
            structure = {'pages': []}
            
            # Extract document metadata
            doc_metadata.update({
                'title': doc.metadata.get('title', ''),
                'author': doc.metadata.get('author', ''),
                'subject': doc.metadata.get('subject', ''),
                'creator': doc.metadata.get('creator', ''),
                'producer': doc.metadata.get('producer', ''),
                'creation_date': doc.metadata.get('creationDate', ''),
                'modification_date': doc.metadata.get('modDate', ''),
                'page_count': len(doc)
            })
            
            # Extract text from each page
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Try to extract text directly
                page_text = page.get_text()
                page_info = {'page_number': page_num + 1, 'has_text': bool(page_text.strip())}
                
                # If no text found and OCR is available, try OCR
                if not page_text.strip() and OCR_AVAILABLE and self.config.get('extract_images', True):
                    notes.append(f"Page {page_num + 1}: Using OCR (no extractable text)")
                    try:
                        # Render page as image
                        pix = page.get_pixmap()
                        img_data = pix.tobytes("png")
                        
                        # OCR the image
                        image = Image.open(io.BytesIO(img_data))
                        page_text = pytesseract.image_to_string(image, lang=self.config['ocr_language'])
                        page_info['ocr_used'] = True
                        
                    except Exception as e:
                        notes.append(f"Page {page_num + 1}: OCR failed - {e}")
                        page_text = f"[OCR FAILED FOR PAGE {page_num + 1}]"
                        page_info['ocr_failed'] = True
                
                if page_text.strip():
                    text_parts.append(f"--- Page {page_num + 1} ---\n{page_text}")
                    page_info['word_count'] = len(page_text.split())
                else:
                    notes.append(f"Page {page_num + 1}: No text content found")
                    page_info['word_count'] = 0
                
                structure['pages'].append(page_info)
            
            doc.close()
            
            return {
                'text': '\n\n'.join(text_parts),
                'metadata': doc_metadata,
                'page_count': len(doc),
                'structure': structure,
                'notes': notes
            }
            
        except Exception as e:
            raise AgentError(f"PDF processing failed: {e}", self.name)
    
    async def _process_docx(self, file_path: Path, metadata: Dict) -> Dict[str, Any]:
        """Process DOCX document with full text extraction"""
        if not DOCX_AVAILABLE:
            raise AgentError("python-docx not available for DOCX processing", self.name)
        
        try:
            doc = DocxDocument(str(file_path))
            
            # Extract text from paragraphs
            text_parts = []
            structure = {'paragraphs': [], 'tables': [], 'headers': []}
            
            # Process paragraphs
            for i, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
                    structure['paragraphs'].append({
                        'index': i,
                        'text_length': len(paragraph.text),
                        'style': paragraph.style.name if paragraph.style else 'Normal'
                    })
            
            # Extract tables if present
            for i, table in enumerate(doc.tables):
                table_text = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_text.append(row_text)
                
                if table_text:
                    # Add table to text
                    table_str = '\n'.join(['\t'.join(row) for row in table_text])
                    text_parts.append(f"\n--- Table {i + 1} ---\n{table_str}\n")
                    
                    structure['tables'].append({
                        'index': i,
                        'rows': len(table_text),
                        'columns': len(table_text[0]) if table_text else 0
                    })
            
            # Extract metadata
            doc_metadata = {
                'title': doc.core_properties.title or '',
                'author': doc.core_properties.author or '',
                'subject': doc.core_properties.subject or '',
                'keywords': doc.core_properties.keywords or '',
                'comments': doc.core_properties.comments or '',
                'created': str(doc.core_properties.created) if doc.core_properties.created else '',
                'modified': str(doc.core_properties.modified) if doc.core_properties.modified else '',
                'paragraph_count': len([p for p in doc.paragraphs if p.text.strip()]),
                'table_count': len(doc.tables)
            }
            
            return {
                'text': '\n\n'.join(text_parts),
                'metadata': doc_metadata,
                'page_count': 1,  # DOCX doesn't have clear page boundaries
                'structure': structure,
                'notes': [f"Extracted {len(structure['paragraphs'])} paragraphs and {len(structure['tables'])} tables"]
            }
            
        except Exception as e:
            raise AgentError(f"DOCX processing failed: {e}", self.name)
    
    async def _process_txt(self, file_path: Path, metadata: Dict) -> Dict[str, Any]:
        """Process plain text document"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            text = None
            encoding_used = None
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text = f.read()
                    encoding_used = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if text is None:
                raise AgentError(f"Could not decode text file with any supported encoding", self.name)
            
            # Analyze structure
            lines = text.split('\n')
            structure = {
                'lines': len(lines),
                'empty_lines': len([line for line in lines if not line.strip()]),
                'encoding': encoding_used
            }
            
            # Basic metadata
            file_stat = file_path.stat()
            doc_metadata = {
                'encoding': encoding_used,
                'file_size': file_stat.st_size,
                'modified': str(file_stat.st_mtime),
                'created': str(file_stat.st_ctime),
                'line_count': len(lines)
            }
            
            return {
                'text': text,
                'metadata': doc_metadata,
                'page_count': 1,
                'structure': structure,
                'notes': [f"Decoded with {encoding_used} encoding, {len(lines)} lines"]
            }
            
        except Exception as e:
            raise AgentError(f"Text processing failed: {e}", self.name)
    
    async def _process_markdown(self, file_path: Path, metadata: Dict) -> Dict[str, Any]:
        """Process Markdown document with structure preservation"""
        try:
            # Read markdown file
            with open(file_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            structure = {'headers': [], 'lists': [], 'tables': [], 'links': []}
            
            if MARKDOWN_AVAILABLE:
                # Convert to HTML for structure analysis
                html = markdown.markdown(md_content, extensions=['tables', 'toc'])
                soup = BeautifulSoup(html, 'html.parser')
                
                # Extract structure
                headers = soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])
                structure['headers'] = [{'level': int(h.name[1]), 'text': h.get_text()} for h in headers]
                
                # Extract tables
                tables = soup.find_all('table')
                for i, table in enumerate(tables):
                    rows = table.find_all('tr')
                    structure['tables'].append({
                        'index': i,
                        'rows': len(rows),
                        'has_header': bool(table.find('th'))
                    })
                
                # Get plain text
                text = soup.get_text()
            else:
                # Basic processing without markdown library
                text = md_content
                # Extract headers manually
                headers = re.findall(r'^(#{1,6})\s+(.+)$', md_content, re.MULTILINE)
                structure['headers'] = [{'level': len(h[0]), 'text': h[1]} for h in headers]
            
            # Metadata
            file_stat = file_path.stat()
            doc_metadata = {
                'format': 'markdown',
                'file_size': file_stat.st_size,
                'modified': str(file_stat.st_mtime),
                'header_count': len(structure['headers']),
                'table_count': len(structure['tables'])
            }
            
            return {
                'text': text,
                'metadata': doc_metadata,
                'page_count': 1,
                'structure': structure,
                'notes': [f"Processed markdown with {len(structure['headers'])} headers"]
            }
            
        except Exception as e:
            raise AgentError(f"Markdown processing failed: {e}", self.name)
    
    async def _process_html(self, file_path: Path, metadata: Dict) -> Dict[str, Any]:
        """Process HTML document"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            if MARKDOWN_AVAILABLE:
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # Extract text
                text = soup.get_text()
                
                # Extract structure
                structure = {
                    'title': soup.title.string if soup.title else '',
                    'headers': [{'level': int(h.name[1]), 'text': h.get_text()} 
                               for h in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])],
                    'links': [{'text': a.get_text(), 'href': a.get('href', '')} 
                             for a in soup.find_all('a')],
                    'images': len(soup.find_all('img'))
                }
                
                # Metadata
                meta_tags = soup.find_all('meta')
                doc_metadata = {
                    'title': structure['title'],
                    'header_count': len(structure['headers']),
                    'link_count': len(structure['links']),
                    'image_count': structure['images']
                }
                
                # Extract meta tag info
                for meta in meta_tags:
                    name = meta.get('name', meta.get('property', ''))
                    content = meta.get('content', '')
                    if name and content:
                        doc_metadata[f'meta_{name}'] = content
            else:
                # Basic text extraction
                text = re.sub(r'<[^>]+>', '', html_content)
                doc_metadata = {'format': 'html'}
                structure = {}
            
            return {
                'text': text,
                'metadata': doc_metadata,
                'page_count': 1,
                'structure': structure,
                'notes': ['Extracted text from HTML']
            }
            
        except Exception as e:
            raise AgentError(f"HTML processing failed: {e}", self.name)
    
    async def _process_rtf(self, file_path: Path, metadata: Dict) -> Dict[str, Any]:
        """Process RTF document (basic implementation)"""
        try:
            # Basic RTF processing - remove RTF codes
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            
            # Remove RTF control codes (very basic)
            text = re.sub(r'\\[a-z]+\d*\s?', '', rtf_content)
            text = re.sub(r'[{}]', '', text)
            text = re.sub(r'\s+', ' ', text).strip()
            
            file_stat = file_path.stat()
            doc_metadata = {
                'format': 'rtf',
                'file_size': file_stat.st_size,
                'modified': str(file_stat.st_mtime)
            }
            
            return {
                'text': text,
                'metadata': doc_metadata,
                'page_count': 1,
                'structure': {},
                'notes': ['Basic RTF processing - install striprtf for better results']
            }
            
        except Exception as e:
            raise AgentError(f"RTF processing failed: {e}", self.name)
    
    # === STRUCTURED DATA HANDLERS ===
    
    async def _process_excel(self, file_path: Path, metadata: Dict) -> Dict[str, Any]:
        """Process Excel file and extract structured data"""
        if not EXCEL_AVAILABLE:
            raise AgentError("pandas/openpyxl not available for Excel processing", self.name)
        
        try:
            # Read all sheets
            xl_file = pd.ExcelFile(str(file_path))
            
            tables = []
            sheets = []
            text_content_parts = []
            database_schema = {}
            
            for sheet_name in xl_file.sheet_names:
                try:
                    df = pd.read_excel(xl_file, sheet_name=sheet_name)
                    
                    # Convert to structured format
                    table_data = {
                        'sheet_name': sheet_name,
                        'columns': df.columns.tolist(),
                        'data': df.values.tolist(),
                        'row_count': len(df),
                        'column_count': len(df.columns),
                        'data_types': {col: str(dtype) for col, dtype in df.dtypes.items()}
                    }
                    
                    tables.append(table_data)
                    sheets.append(sheet_name)
                    
                    # Generate database schema suggestion
                    schema = {
                        'table_name': f"{file_path.stem}_{sheet_name}".lower().replace(' ', '_'),
                        'columns': []
                    }
                    
                    for col, dtype in df.dtypes.items():
                        sql_type = self._pandas_to_sql_type(dtype)
                        schema['columns'].append({
                            'name': str(col).lower().replace(' ', '_'),
                            'type': sql_type,
                            'original_name': str(col)
                        })
                    
                    database_schema[sheet_name] = schema
                    
                    # Extract any text content for search
                    text_values = df.select_dtypes(include=['object']).values.flatten()
                    text_content = ' '.join([str(val) for val in text_values if pd.notna(val)])
                    if text_content.strip():
                        text_content_parts.append(f"Sheet: {sheet_name}\n{text_content}")
                
                except Exception as e:
                    logger.warning(f"Could not process sheet '{sheet_name}': {e}")
            
            # Metadata
            doc_metadata = {
                'format': 'excel',
                'sheet_count': len(xl_file.sheet_names),
                'sheet_names': xl_file.sheet_names,
                'total_tables': len(tables),
                'file_size': file_path.stat().st_size
            }
            
            return {
                'tables': tables,
                'sheets': sheets,
                'database_schema': database_schema,
                'text_content': '\n\n'.join(text_content_parts),
                'metadata': doc_metadata,
                'notes': [f"Extracted {len(tables)} sheets with structured data"]
            }
            
        except Exception as e:
            raise AgentError(f"Excel processing failed: {e}", self.name)
    
    async def _process_csv(self, file_path: Path, metadata: Dict) -> Dict[str, Any]:
        """Process CSV file and extract structured data"""
        if not EXCEL_AVAILABLE:
            raise AgentError("pandas not available for CSV processing", self.name)
        
        try:
            # Try different encodings and delimiters
            encodings = ['utf-8', 'latin-1', 'cp1252']
            delimiters = [',', ';', '\t', '|']
            
            df = None
            encoding_used = None
            delimiter_used = None
            
            for encoding in encodings:
                for delimiter in delimiters:
                    try:
                        df = pd.read_csv(file_path, encoding=encoding, sep=delimiter)
                        if len(df.columns) > 1:  # Valid CSV should have multiple columns
                            encoding_used = encoding
                            delimiter_used = delimiter
                            break
                    except:
                        continue
                if df is not None:
                    break
            
            if df is None:
                raise AgentError("Could not parse CSV with any supported encoding/delimiter", self.name)
            
            # Structure data
            table_data = {
                'sheet_name': 'data',
                'columns': df.columns.tolist(),
                'data': df.values.tolist(),
                'row_count': len(df),
                'column_count': len(df.columns),
                'data_types': {col: str(dtype) for col, dtype in df.dtypes.items()}
            }
            
            # Generate database schema
            schema = {
                'table_name': file_path.stem.lower().replace(' ', '_'),
                'columns': []
            }
            
            for col, dtype in df.dtypes.items():
                sql_type = self._pandas_to_sql_type(dtype)
                schema['columns'].append({
                    'name': str(col).lower().replace(' ', '_'),
                    'type': sql_type,
                    'original_name': str(col)
                })
            
            # Extract text content
            text_values = df.select_dtypes(include=['object']).values.flatten()
            text_content = ' '.join([str(val) for val in text_values if pd.notna(val)])
            
            # Metadata
            doc_metadata = {
                'format': 'csv',
                'encoding': encoding_used,
                'delimiter': delimiter_used,
                'row_count': len(df),
                'column_count': len(df.columns),
                'file_size': file_path.stat().st_size
            }
            
            return {
                'tables': [table_data],
                'sheets': ['data'],
                'database_schema': {'data': schema},
                'text_content': text_content,
                'metadata': doc_metadata,
                'notes': [f"Parsed CSV with {delimiter_used} delimiter and {encoding_used} encoding"]
            }
            
        except Exception as e:
            raise AgentError(f"CSV processing failed: {e}", self.name)
    
    # === REFERENCE ONLY HANDLERS ===
    
    async def _process_powerpoint(self, file_path: Path, metadata: Dict) -> Dict[str, Any]:
        """Process PowerPoint file for reference (metadata + basic text)"""
        if not PPTX_AVAILABLE:
            # Fallback to basic file info
            return self._basic_file_reference(file_path, 'powerpoint')
        
        try:
            prs = Presentation(str(file_path))
            
            text_parts = []
            slide_info = []
            
            # Extract basic text from slides
            for i, slide in enumerate(prs.slides):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                slide_content = ' '.join(slide_text)
                if slide_content:
                    text_parts.append(f"Slide {i + 1}: {slide_content}")
                
                slide_info.append({
                    'slide_number': i + 1,
                    'text_length': len(slide_content),
                    'has_content': bool(slide_content.strip())
                })
            
            preview_text = '\n'.join(text_parts[:3])  # First 3 slides for preview
            
            doc_metadata = {
                'format': 'powerpoint',
                'slide_count': len(prs.slides),
                'slides_with_content': len([s for s in slide_info if s['has_content']]),
                'file_size': file_path.stat().st_size
            }
            
            return {
                'preview_text': preview_text,
                'slide_count': len(prs.slides),
                'metadata': doc_metadata,
                'summary': f"PowerPoint presentation with {len(prs.slides)} slides",
                'notes': [f"Extracted basic text from {len(prs.slides)} slides"]
            }
            
        except Exception as e:
            logger.warning(f"PowerPoint processing failed: {e}")
            return self._basic_file_reference(file_path, 'powerpoint')
    
    async def _process_image_reference(self, file_path: Path, metadata: Dict) -> Dict[str, Any]:
        """Process image file for reference (with optional OCR preview)"""
        try:
            preview_text = ""
            image_info = {}
            
            if OCR_AVAILABLE:
                try:
                    # Quick OCR for preview (first 200 characters)
                    image = Image.open(file_path)
                    full_text = pytesseract.image_to_string(image, lang=self.config['ocr_language'])
                    preview_text = full_text[:200] + "..." if len(full_text) > 200 else full_text
                    
                    image_info = {
                        'format': image.format,
                        'mode': image.mode,
                        'size': image.size,
                        'ocr_available': True
                    }
                except Exception as e:
                    logger.debug(f"OCR preview failed for {file_path.name}: {e}")
                    image_info = {'ocr_available': False}
            
            doc_metadata = {
                'format': 'image',
                'file_size': file_path.stat().st_size,
                **image_info
            }
            
            return {
                'preview_text': preview_text,
                'image_count': 1,
                'metadata': doc_metadata,
                'summary': f"Image file ({file_path.suffix.upper()}) - OCR available",
                'notes': ['Image processed for reference with OCR preview']
            }
            
        except Exception as e:
            return self._basic_file_reference(file_path, 'image')
    
    def _basic_file_reference(self, file_path: Path, file_type: str) -> Dict[str, Any]:
        """Create basic file reference when full processing fails"""
        file_stat = file_path.stat()
        
        return {
            'preview_text': f"File reference: {file_path.name}",
            'metadata': {
                'format': file_type,
                'file_size': file_stat.st_size,
                'modified': str(file_stat.st_mtime),
                'accessible': True
            },
            'summary': f"{file_type.title()} file available for reference",
            'notes': [f"Basic file reference - processing library not available"]
        }
    
    # === UTILITY METHODS ===
    
    def _pandas_to_sql_type(self, pandas_dtype) -> str:
        """Convert pandas dtype to SQL type"""
        dtype_str = str(pandas_dtype)
        
        if 'int' in dtype_str:
            return 'INTEGER'
        elif 'float' in dtype_str:
            return 'REAL'
        elif 'bool' in dtype_str:
            return 'BOOLEAN'
        elif 'datetime' in dtype_str:
            return 'TIMESTAMP'
        else:
            return 'TEXT'
    
    def _calculate_content_hash(self, text: str) -> str:
        """Calculate hash of document content for duplicate detection"""
        normalized = ' '.join(text.lower().split())
        return hashlib.sha256(normalized.encode('utf-8')).hexdigest()
    
    def _is_legal_document(self, text: str) -> bool:
        """Determine if document appears to be legal in nature"""
        if not text or len(text) < 100:  # Too short to determine
            return False
        
        text_lower = text.lower()
        indicator_count = sum(1 for indicator in self.legal_indicators if indicator in text_lower)
        
        # Consider it legal if it has multiple legal terms
        return indicator_count >= 3
    
    def _classify_legal_document(self, text: str) -> Dict[str, Any]:
        """Classify the type of legal document"""
        text_lower = text.lower()
        
        classifications = {
            'motion': ['motion to', 'motion for', 'motion that'],
            'complaint': ['complaint', 'plaintiff', 'defendant', 'cause of action'],
            'affidavit': ['affidavit', 'affiant', 'sworn statement', 'under oath'],
            'deposition': ['deposition', 'examination under oath', 'q:', 'a:'],
            'court_order': ['order', 'it is ordered', 'court orders', 'hereby ordered'],
            'warrant': ['warrant', 'search warrant', 'arrest warrant', 'probable cause'],
            'brief': ['brief', 'argument', 'legal memorandum', 'statement of facts'],
            'contract': ['agreement', 'contract', 'party of the first part', 'whereas'],
            'statute': ['statute', 'section', 'subsection', 'chapter', 'code'],
            'constitution': ['constitution', 'amendment', 'article', 'bill of rights'],
            'evidence_log': ['evidence', 'exhibit', 'chain of custody', 'collected'],
            'witness_statement': ['witness', 'statement', 'testified', 'observed']
        }
        
        detected_types = []
        confidence_scores = {}
        
        for doc_type, indicators in classifications.items():
            score = sum(1 for indicator in indicators if indicator in text_lower)
            if score > 0:
                detected_types.append(doc_type)
                confidence_scores[doc_type] = score
        
        primary_type = max(confidence_scores.items(), key=lambda x: x[1])[0] if confidence_scores else 'unknown'
        
        return {
            'primary_type': primary_type,
            'detected_types': detected_types,
            'confidence_scores': confidence_scores
        }
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text"""
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove page numbers and headers/footers patterns
        text = re.sub(r'\n\s*\d+\s*\n', '\n', text)
        text = re.sub(r'\n\s*Page \d+ of \d+\s*\n', '\n', text, flags=re.IGNORECASE)
        
        # Remove excessive line breaks
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        # Clean up common OCR artifacts
        text = re.sub(r'[^\x00-\x7F]+', ' ', text)  # Remove non-ASCII
        text = re.sub(r'(\w)\1{3,}', r'\1\1', text)  # Reduce repeated characters
        
        return text.strip()
    
    # === PUBLIC INTERFACE METHODS ===
    
    def get_supported_formats(self) -> Dict[str, Dict[str, str]]:
        """Get list of supported file formats with their processing strategies"""
        return {ext: {'type': info['type'], 'strategy': info['strategy']} 
                for ext, info in self.file_config.items()}
    
    def get_processing_strategy(self, file_path: Union[str, Path]) -> str:
        """Get the processing strategy for a file"""
        file_info = self._detect_document_info(Path(file_path))
        return file_info['strategy']
    
    async def process_directory(self, directory_path: Union[str, Path],
                               recursive: bool = False,
                               file_patterns: List[str] = None,
                               strategy_filter: str = None) -> List[AgentResult]:
        """
        Process all supported files in a directory with strategy filtering
        
        Args:
            directory_path: Path to directory
            recursive: Whether to search subdirectories
            file_patterns: List of file patterns to match
            strategy_filter: Only process files with this strategy (e.g., 'full_processing')
            
        Returns:
            List of AgentResults for each processed file
        """
        directory = Path(directory_path)
        if not directory.exists() or not directory.is_dir():
            raise AgentError(f"Directory not found: {directory}", self.name)
        
        # Find files to process
        files_to_process = []
        
        search_pattern = '**/*' if recursive else '*'
        
        for file_path in directory.glob(search_pattern):
            if file_path.is_file():
                file_info = self._detect_document_info(file_path)
                
                # Check if file type is supported
                if file_info['type'] != DocumentType.UNKNOWN:
                    # Apply strategy filter if specified
                    if strategy_filter and file_info['strategy'] != strategy_filter:
                        continue
                    
                    # Apply file patterns if specified
                    if file_patterns:
                        if any(file_path.match(pattern) for pattern in file_patterns):
                            files_to_process.append(file_path)
                    else:
                        files_to_process.append(file_path)
        
        logger.info(f"Found {len(files_to_process)} files to process in {directory}")
        
        # Process files concurrently (with limit)
        semaphore = asyncio.Semaphore(3)
        
        async def process_with_semaphore(file_path):
            async with semaphore:
                return await self.process(file_path)
        
        # Process all files
        results = await asyncio.gather(
            *[process_with_semaphore(file_path) for file_path in files_to_process],
            return_exceptions=True
        )
        
        # Convert exceptions to failed results
        processed_results = []
        for i, result in enumerate(results):
            if isinstance(result, Exception):
                processed_results.append(AgentResult(
                    success=False,
                    error=str(result),
                    metadata={'file_path': str(files_to_process[i])},
                    agent_name=self.name
                ))
            else:
                processed_results.append(result)
        
        return processed_results