# legal_ai_system/agents/document_processor/document_processor_agent.py
"""
Consolidated Document Processor Agent.
Processes various file formats (PDF, DOCX, TXT, MD, HTML, RTF, Excel, CSV, PPTX, Images)
for text extraction, metadata retrieval, and structural analysis.
Integrates optional dependencies gracefully and uses shared components.
"""


from __future__ import annotations

import asyncio
import hashlib
import io
import mimetypes
import json
import zipfile
import email
import re
from dataclasses import asdict, dataclass, field
from datetime import datetime, timezone
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

from ..core.agent_unified_config import create_agent_memory_mixin
from ..core.base_agent import BaseAgent
from ..core.constants import Constants
from ..core.detailed_logging import (
    LogCategory,
    detailed_log_function,
    get_detailed_logger,
)
from ..core.unified_exceptions import AgentExecutionError, DocumentProcessingError
from ..utils.dependency_manager import DependencyManager

MemoryMixin = create_agent_memory_mixin()
file_logger = get_detailed_logger("File_Processing", LogCategory.DOCUMENT)
# Attempt to import optional dependencies via DependencyManager
dep_manager = DependencyManager()
dep_manager.check_dependencies()

# Conditional imports based on availability (global for the module)
if dep_manager.is_available("pymupdf"):
    import fitz  # PyMuPDF
if dep_manager.is_available("docx"):
    from docx import Document as DocxDocument
    from docx.opc.exceptions import PackageNotFoundError as DocxPackageNotFoundError
if dep_manager.is_available("pytesseract") and dep_manager.is_available("PIL"):
    import pytesseract
    from PIL import Image, UnidentifiedImageError
if dep_manager.is_available("pandas"):
    import pandas as pd
if dep_manager.is_available("pptx"):
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError
if dep_manager.is_available("markdown") and dep_manager.is_available("bs4"):
    import markdown
    from bs4 import BeautifulSoup
if dep_manager.is_available("striprtf"):  # For RTF
    from striprtf.striprtf import rtf_to_text


class ProcessingStrategy:
    FULL_PROCESSING = "full_processing"  # Extract all text, metadata, structure
    STRUCTURED_DATA = "structured_data"  # Primarily for Excel, CSV - extract tables
    REFERENCE_ONLY = "reference_only"  # Extract metadata, maybe slide titles (PPTX)
    OCR_IF_NEEDED = "ocr_if_needed"  # For images, or PDFs with no text layer
    METADATA_ONLY = "metadata_only"  # Only extract file system and basic metadata


class DocumentContentType(Enum):
    PDF = "application/pdf"
    DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    DOC = "application/msword"
    TXT = "text/plain"
    MD = "text/markdown"
    HTML = "text/html"
    RTF = "application/rtf"
    XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    XLS = "application/vnd.ms-excel"
    CSV = "text/csv"
    PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    PPT = "application/vnd.ms-powerpoint"
    PNG = "image/png"
    JPG = "image/jpeg"
    TIFF = "image/tiff"
    BMP = "image/bmp"
    MP3 = "audio/mpeg"
    WAV = "audio/wav"
    MP4 = "video/mp4"
    MOV = "video/quicktime"
    EML = "message/rfc822"
    ZIP = "application/zip"
    JSON = "application/json"
    YAML = "application/x-yaml"
    UNKNOWN = "application/octet-stream"

    @classmethod
    def from_extension(cls, ext: str) -> "DocumentContentType":
        ext_map = {
            ".pdf": cls.PDF,
            ".docx": cls.DOCX,
            ".doc": cls.DOC,
            ".txt": cls.TXT,
            ".md": cls.MD,
            ".html": cls.HTML,
            ".htm": cls.HTML,
            ".rtf": cls.RTF,
            ".xlsx": cls.XLSX,
            ".xls": cls.XLS,
            ".csv": cls.CSV,
            ".pptx": cls.PPTX,
            ".ppt": cls.PPT,
            ".png": cls.PNG,
            ".jpg": cls.JPG,
            ".jpeg": cls.JPG,
            ".tiff": cls.TIFF,
            ".tif": cls.TIFF,
            ".bmp": cls.BMP,
            ".mp3": cls.MP3,
            ".wav": cls.WAV,
            ".mp4": cls.MP4,
            ".mov": cls.MOV,
            ".eml": cls.EML,
            ".zip": cls.ZIP,
            ".json": cls.JSON,
            ".yaml": cls.YAML,
            ".yml": cls.YAML,
        }
        return ext_map.get(ext.lower(), cls.UNKNOWN)

    @classmethod
    def from_mimetype(cls, mime: Optional[str]) -> "DocumentContentType":
        if not mime:
            return cls.UNKNOWN
        for member in cls:
            if member.value == mime:
                return member
        # Fallback for common variations
        if "word" in mime:
            return cls.DOCX  # Prefer modern if unsure
        if "excel" in mime or "spreadsheet" in mime:
            return cls.XLSX
        if "presentation" in mime:
            return cls.PPTX
        if "rfc822" in mime or "message" in mime:
            return cls.EML
        if "zip" in mime:
            return cls.ZIP
        if "json" in mime:
            return cls.JSON
        if "yaml" in mime or "yml" in mime:
            return cls.YAML
        if "audio" in mime:
            if "mpeg" in mime:
                return cls.MP3
            if "wav" in mime:
                return cls.WAV
        if "video" in mime:
            if "mp4" in mime:
                return cls.MP4
            if "quicktime" in mime:
                return cls.MOV
        return cls.UNKNOWN


@dataclass
class DocumentProcessingOutput:
    file_path: str
    file_name: str
    file_size_bytes: int
    document_content_type: DocumentContentType
    processing_strategy_used: str

    text_content: Optional[str] = None
    content_hash_sha256: Optional[str] = None  # Renamed for clarity
    is_legal_document_classified: Optional[bool] = None
    classification_details: Optional[Dict[str, Any]] = None

    extracted_metadata: Dict[str, Any] = field(
        default_factory=dict
    )  # From file itself + file system
    page_count: Optional[int] = None
    word_count: Optional[int] = None
    character_count: Optional[int] = None
    document_structure: Optional[Dict[str, Any]] = None  # e.g., headings, sections

    tables_extracted: Optional[List[Dict[str, Any]]] = (
        None  # List of tables, each table a list of lists or dict
    )
    table_count: Optional[int] = None

    image_ocr_results: Optional[List[Dict[str, str]]] = (
        None  # For files with embedded images + OCR
    )

    processing_notes: List[str] = field(default_factory=list)
    errors: List[str] = field(default_factory=list)
    processing_time_sec: float = 0.0
    processed_at: str = field(
        default_factory=lambda: datetime.now(timezone.utc).isoformat()
    )

    def to_dict(self) -> Dict[str, Any]:
        data = asdict(self)
        data["document_content_type"] = (
            self.document_content_type.value
        )  # Store enum value
        return data


class DocumentProcessorAgent(BaseAgent, MemoryMixin):
    """
    Consolidated agent for processing various legal document formats.
    Extracts text, metadata, and basic structure. Handles optional dependencies.
    """

    def __init__(self, service_container: Any, **config: Any):
        super().__init__(
            service_container,
            name="DocumentProcessorAgent",
            agent_type="document_processing",
        )
        # Merge any passed configuration overrides
        if config:
            self.config.update(config)

        # Get optimized Grok-Mini configuration for this agent
        self.llm_config = self.get_optimized_llm_config()
        self.logger.info(
            f"DocumentProcessorAgent configured with model: {self.llm_config.get('llm_model', 'default')}"
        )
        # Shared components (can be injected or fetched)
        # NLP classifier uses ML model when available and falls back to keywords
        self.classifier = NLPDocumentClassifier()
        self.chunker = DocumentChunker(
            chunk_size=self.config.get("dp_chunk_size", 4000),
            overlap=self.config.get("dp_chunk_overlap", 400),
        )

        # Configuration
        self.max_file_size_mb = float(
            self.config.get("max_file_size_mb", Constants.Document.MAX_DOCUMENT_SIZE_MB)
        )
        self.ocr_language = str(self.config.get("ocr_language", "eng"))
        self.ocr_dpi = int(self.config.get("ocr_dpi", 300))
        self.extract_images_from_pdf = bool(
            self.config.get("extract_images_from_pdf", False)
        )  # Usually not needed for text processing
        self.clean_extracted_text = bool(self.config.get("clean_extracted_text", True))
        self.preserve_document_structure = bool(
            self.config.get("preserve_document_structure", False)
        )  # Default to False if not fully supported
        self.extract_tables_from_docs = bool(
            self.config.get("extract_tables_from_docs", True)
        )

        self._initialize_file_type_configs()

        self.logger.info(
            f"{self.name} initialized.", parameters=self.get_config_summary_params()
        )

    def _initialize_file_type_configs(self):
        """Defines processing handlers and strategies for different file types."""
        self.file_type_configs: Dict[DocumentContentType, Dict[str, Any]] = {
            DocumentContentType.PDF: {
                "strategy": ProcessingStrategy.FULL_PROCESSING,
                "handler_method": "_process_pdf_async",
                "deps": ["pymupdf"],
            },
            DocumentContentType.DOCX: {
                "strategy": ProcessingStrategy.FULL_PROCESSING,
                "handler_method": "_process_docx_async",
                "deps": ["docx"],
            },
            DocumentContentType.DOC: {
                "strategy": ProcessingStrategy.FULL_PROCESSING,
                "handler_method": "_process_doc_async",
                "deps": ["docx"],
            },  # Often python-docx can handle .doc via Word interop or conversion if available
            DocumentContentType.TXT: {
                "strategy": ProcessingStrategy.FULL_PROCESSING,
                "handler_method": "_process_txt_async",
                "deps": [],
            },
            DocumentContentType.MD: {
                "strategy": ProcessingStrategy.FULL_PROCESSING,
                "handler_method": "_process_markdown_async",
                "deps": ["markdown", "bs4"],
            },
            DocumentContentType.HTML: {
                "strategy": ProcessingStrategy.FULL_PROCESSING,
                "handler_method": "_process_html_async",
                "deps": ["bs4"],
            },
            DocumentContentType.RTF: {
                "strategy": ProcessingStrategy.FULL_PROCESSING,
                "handler_method": "_process_rtf_async",
                "deps": ["striprtf"],
            },
            DocumentContentType.XLSX: {
                "strategy": ProcessingStrategy.STRUCTURED_DATA,
                "handler_method": "_process_excel_async",
                "deps": ["pandas", "openpyxl"],
            },
            DocumentContentType.XLS: {
                "strategy": ProcessingStrategy.STRUCTURED_DATA,
                "handler_method": "_process_excel_async",
                "deps": ["pandas", "xlrd"],
            },  # xlrd for older .xls
            DocumentContentType.CSV: {
                "strategy": ProcessingStrategy.STRUCTURED_DATA,
                "handler_method": "_process_csv_async",
                "deps": ["pandas"],
            },
            DocumentContentType.PPTX: {
                "strategy": ProcessingStrategy.REFERENCE_ONLY,
                "handler_method": "_process_powerpoint_async",
                "deps": ["pptx"],
            },
            DocumentContentType.PPT: {
                "strategy": ProcessingStrategy.REFERENCE_ONLY,
                "handler_method": "_process_powerpoint_async",
                "deps": ["pptx"],
            },  # python-pptx does not handle .ppt directly
            DocumentContentType.PNG: {
                "strategy": ProcessingStrategy.OCR_IF_NEEDED,
                "handler_method": "_process_image_async",
                "deps": ["pytesseract", "PIL"],
            },
            DocumentContentType.JPG: {
                "strategy": ProcessingStrategy.OCR_IF_NEEDED,
                "handler_method": "_process_image_async",
                "deps": ["pytesseract", "PIL"],
            },
            DocumentContentType.TIFF: {
                "strategy": ProcessingStrategy.OCR_IF_NEEDED,
                "handler_method": "_process_image_async",
                "deps": ["pytesseract", "PIL"],
            },
            DocumentContentType.BMP: {
                "strategy": ProcessingStrategy.OCR_IF_NEEDED,
                "handler_method": "_process_image_async",
                "deps": ["pytesseract", "PIL"],
            },
            DocumentContentType.MP3: {
                "strategy": ProcessingStrategy.METADATA_ONLY,
                "handler_method": "_process_audio_async",
                "deps": ["pydub"],
            },
            DocumentContentType.WAV: {
                "strategy": ProcessingStrategy.METADATA_ONLY,
                "handler_method": "_process_audio_async",
                "deps": ["pydub"],
            },
            DocumentContentType.MP4: {
                "strategy": ProcessingStrategy.METADATA_ONLY,
                "handler_method": "_process_video_async",
                "deps": ["moviepy"],
            },
            DocumentContentType.MOV: {
                "strategy": ProcessingStrategy.METADATA_ONLY,
                "handler_method": "_process_video_async",
                "deps": ["moviepy"],
            },
            DocumentContentType.EML: {
                "strategy": ProcessingStrategy.FULL_PROCESSING,
                "handler_method": "_process_eml_async",
                "deps": [],
            },
            DocumentContentType.ZIP: {
                "strategy": ProcessingStrategy.METADATA_ONLY,
                "handler_method": "_process_zip_async",
                "deps": [],
            },
            DocumentContentType.JSON: {
                "strategy": ProcessingStrategy.FULL_PROCESSING,
                "handler_method": "_process_json_async",
                "deps": [],
            },
            DocumentContentType.YAML: {
                "strategy": ProcessingStrategy.FULL_PROCESSING,
                "handler_method": "_process_yaml_async",
                "deps": ["yaml"],
            },
        }

    def get_config_summary_params(self) -> Dict[str, Any]:
        return {
            "max_file_size_mb": self.max_file_size_mb,
            "ocr_lang": self.ocr_language,
            "ocr_dpi": self.ocr_dpi,
            "clean_text": self.clean_extracted_text,
            "preserve_structure": self.preserve_document_structure,
            "extract_tables": self.extract_tables_from_docs,
        }

    @detailed_log_function(LogCategory.DOCUMENT)  # Use agent's logger with category
    async def _process_task(
        self, task_data: Union[str, Path], metadata: Dict[str, Any]
    ) -> Dict[str, Any]:
        start_time = datetime.now(timezone.utc)
        file_path = Path(task_data)  # Assume task_data is the file path
        self.logger.info(f"Starting document processing for: {file_path.name}")

        # Initialize output with basic file info
        try:
            output = DocumentProcessingOutput(
                file_path=str(file_path),
                file_name=file_path.name,
                file_size_bytes=file_path.stat().st_size if file_path.exists() else 0,
                document_content_type=DocumentContentType.UNKNOWN,  # Will be updated
                processing_strategy_used="unknown",
            )
        except FileNotFoundError:
            self.logger.error(f"File not found at path: {file_path}")
            # Return a minimal error output directly
            error_output = DocumentProcessingOutput(
                file_path=str(file_path),
                file_name=file_path.name,
                file_size_bytes=0,
                document_content_type=DocumentContentType.UNKNOWN,
                processing_strategy_used="error",
            )
            error_output.errors.append(f"File not found: {file_path}")
            error_output.processing_time_sec = round(
                (datetime.now(timezone.utc) - start_time).total_seconds(), 3
            )
            return error_output.to_dict()

        try:
            if not file_path.is_file():
                raise DocumentProcessingError(
                    f"Path is not a file: {file_path}", file_path=file_path
                )

            if (output.file_size_bytes / (1024 * 1024)) > self.max_file_size_mb:
                raise DocumentProcessingError(
                    f"File '{file_path.name}' ({output.file_size_bytes / (1024*1024):.2f}MB) "
                    f"exceeds max size of {self.max_file_size_mb}MB.",
                    file_path=file_path,
                )

            # Determine file type and handler
            file_ext = file_path.suffix.lower()
            output.document_content_type = DocumentContentType.from_extension(file_ext)

            # If extension not recognized, try MIME type
            if output.document_content_type == DocumentContentType.UNKNOWN:
                mime_type_str, _ = mimetypes.guess_type(str(file_path))
                output.extracted_metadata["mime_type_detected"] = mime_type_str
                output.document_content_type = DocumentContentType.from_mimetype(
                    mime_type_str
                )
                self.logger.debug(
                    f"Used MIME type '{mime_type_str}' to determine content type: {output.document_content_type.name}",
                    parameters={"file": file_path.name},
                )

            file_type_enum = output.document_content_type
            file_handler_config = self.file_type_configs.get(file_type_enum)

            if not file_handler_config:
                raise DocumentProcessingError(
                    f"Unsupported file type: {file_type_enum.name} (ext: {file_ext})",
                    file_path=file_path,
                )

            output.processing_strategy_used = file_handler_config["strategy"]

            # Check dependencies
            for dep_key in file_handler_config.get("deps", []):
                if not dep_manager.is_available(dep_key):
                    dep_info = dep_manager.DEPENDENCY_MAP.get(
                        dep_key, {"package": f"unknown_package_for_{dep_key}"}
                    )
                    raise DocumentProcessingError(
                        f"Missing dependency '{dep_info['package']}' required for {file_type_enum.name} files.",
                        file_path=file_path,
                    )

            handler_method_name = file_handler_config["handler_method"]
            handler_method = getattr(self, handler_method_name, None)
            if not handler_method or not callable(handler_method):
                raise DocumentProcessingError(
                    f"Internal error: Handler method '{handler_method_name}' not found or not callable.",
                    file_path=file_path,
                )

            # Pass relevant processing options to the handler method
            processing_options = {
                "force_ocr": metadata.get(
                    "force_ocr", False
                ),  # Metadata can override agent defaults
                "ocr_language": self.ocr_language,
                "ocr_dpi": self.ocr_dpi,
                "extract_images": self.extract_images_from_pdf,
                "clean_text": self.clean_extracted_text,
                "preserve_structure": self.preserve_document_structure,
                "extract_tables": self.extract_tables_from_docs,
                "target_chunk_size": self.chunker.chunk_size,  # For handlers that might pre-chunk
                "metadata_extraction_level": metadata.get(
                    "metadata_extraction_level", "full"
                ),  # 'basic' or 'full'
            }

            handler_result = await asyncio.to_thread(
                handler_method, file_path, processing_options
            )

            # Populate output from handler result
            output.text_content = handler_result.get("text_content")
            output.extracted_metadata.update(
                handler_result.get("extracted_metadata", {})
            )
            output.page_count = handler_result.get("page_count")
            output.document_structure = handler_result.get("document_structure")
            output.tables_extracted = handler_result.get("tables_extracted")
            output.table_count = (
                len(output.tables_extracted)
                if output.tables_extracted
                else handler_result.get("table_count", 0)
            )
            output.image_ocr_results = handler_result.get("image_ocr_results")
            output.processing_notes.extend(handler_result.get("processing_notes", []))

            if output.text_content:
                output.word_count = len(output.text_content.split())
                output.character_count = len(output.text_content)
                # Use a faster hash if performance is critical for many small files, SHA256 is robust.
                output.content_hash_sha256 = hashlib.sha256(
                    output.text_content.encode("utf-8", "ignore")
                ).hexdigest()

                # Perform classification using the shared NLPDocumentClassifier
                classification_result = self.classifier.classify(
                    output.text_content, filename=file_path.name
                )
                output.is_legal_document_classified = classification_result.get(
                    "is_legal_document", False
                )
                output.classification_details = classification_result
                output.processing_notes.append(
                    f"Document classification: Type='{classification_result.get('primary_type', 'N/A')}', Score={classification_result.get('primary_score', 0.0):.2f}"
                )

            self.logger.info(
                f"Successfully processed document: {file_path.name}",
                parameters={
                    "type": output.document_content_type.name,
                    "text_len": output.character_count or 0,
                },
            )

        except DocumentProcessingError as dpe:
            self.logger.error(
                f"DocumentProcessingError for '{file_path.name}': {dpe.message}",
                parameters={
                    "file": str(file_path),
                    "error_details": dpe.technical_details,
                },
            )
            output.errors.append(dpe.message)
        except AgentExecutionError as aee:
            self.logger.error(
                f"AgentExecutionError during processing of '{file_path.name}': {str(aee)}",
                parameters={"file": str(file_path)},
                exception=aee,
            )  # Log full exception if needed
            output.errors.append(f"Agent execution error: {str(aee)}")
        except Exception as e:
            self.logger.error(
                f"Unexpected error processing '{file_path.name}'.",
                parameters={"file": str(file_path)},
                exception=e,
                exc_info=True,
            )  # exc_info for unexpected
            output.errors.append(
                f"Unexpected critical error: {type(e).__name__} - {str(e)}"
            )

        finally:
            output.processing_time_sec = round(
                (datetime.now(timezone.utc) - start_time).total_seconds(), 3
            )
            if output.errors:
                self.logger.warning(
                    f"Processing for '{file_path.name}' completed with {len(output.errors)} error(s).",
                    parameters={
                        "first_error": output.errors[0] if output.errors else "N/A"
                    },
                )

        return output.to_dict()

    # --- ASYNC WRAPPERS FOR HANDLERS ---

    # --- SYNCHRONOUS FILE PROCESSING LOGIC ---
    def _common_text_cleaning(self, text: Optional[str]) -> Optional[str]:
        if not text:
            return None
        text = re.sub(r"\s+", " ", text)  # Replace multiple whitespaces with single
        text = re.sub(r"\n{3,}", "\n\n", text)  # Reduce multiple newlines to max two
        text = re.sub(
            r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]", "", text
        )  # Remove most control characters except \t, \n, \r
        return text.strip()

    def _sync_process_pdf(
        self, file_path: Path, options: Dict[str, Any]
    ) -> Dict[str, Any]:
        file_logger.info(f"Processing PDF (sync): {file_path.name}")
        if not dep_manager.is_available("pymupdf"):
            raise DocumentProcessingError(
                "PyMuPDF (fitz) unavailable.", file_path=file_path
            )

        text_content_parts: List[str] = []
        extracted_meta: Dict[str, Any] = {"source_format": "pdf"}
        notes: List[str] = []
        page_count = 0
        image_ocr_res: List[Dict[str, str]] = []
        file_hash = ocr_cache.compute_file_hash(file_path)

        try:
            doc = fitz.open(str(file_path))  # type: ignore
            page_count = len(doc)
            meta = doc.metadata or {}
            for k, v in meta.items():  # Clean metadata values
                if isinstance(v, str) and v.startswith("D:"):
                    meta[k] = v[2:]  # Clean PDF date format
            extracted_meta.update(meta)
            extracted_meta["page_count_from_meta"] = (
                page_count  # Store original page count from metadata
            )

            for i in range(page_count):
                page = doc.load_page(i)
                cached = ocr_cache.get(file_path, page=i + 1, file_hash=file_hash)
                if cached and not options.get("force_ocr", False):
                    page_text = cached
                    notes.append(f"Page {i+1}: Loaded OCR text from cache.")
                else:
                    page_text = page.get_text(  # type: ignore[attr-defined]
                        "text", sort=True
                    ).strip()
                    needs_ocr = options.get("force_ocr", False) or (
                        not page_text and options.get("ocr_if_needed", True)
                    )
                    if needs_ocr:
                        if dep_manager.is_available("pytesseract") and dep_manager.is_available("PIL"):
                            notes.append(
                                f"Page {i+1}: Attempting OCR (Reason: {'forced' if options.get('force_ocr') else 'no text layer'})."
                            )
                            try:
                                pix = page.get_pixmap(  # type: ignore[attr-defined]
                                    dpi=options.get("ocr_dpi", 300)
                                )
                                img_bytes = pix.tobytes(  # type: ignore[attr-defined]
                                    "png"
                                )
                                pil_img = Image.open(io.BytesIO(img_bytes))  # type: ignore
                                ocr_text = pytesseract.image_to_string(
                                    pil_img,
                                    lang=options.get("ocr_language", "eng"),
                                ).strip()  # type: ignore
                                if ocr_text:
                                    page_text = ocr_text
                                    ocr_cache.set(
                                        file_path,
                                        page=i + 1,
                                        text=ocr_text,
                                        file_hash=file_hash,
                                    )
                                    notes.append(
                                        f"Page {i+1}: OCR successful, added {len(ocr_text)} chars."
                                    )
                                    image_ocr_res.append(
                                        {
                                            "page": str(i + 1),
                                            "ocr_text_preview": ocr_text[:100] + "...",
                                        }
                                    )
                                else:
                                    notes.append(f"Page {i+1}: OCR yielded no text.")
                            except Exception as ocr_e:
                                notes.append(
                                    f"Page {i+1}: OCR failed - {str(ocr_e)[:100]}."
                                )
                                file_logger.warning(
                                    f"OCR for page {i+1} of {file_path.name} failed.",
                                    exception=ocr_e,
                                )
                        else:
                            notes.append(
                                f"Page {i+1}: OCR skipped (pytesseract/Pillow unavailable)."
                            )

                if page_text:
                    text_content_parts.append(page_text)
            doc.close()
        except Exception as e:
            file_logger.error(f"Error processing PDF {file_path.name}.", exception=e)
            notes.append(f"Critical PDF processing error: {str(e)}")
            # For PDF, if extraction fails, there might be no text at all.
            # Consider if DocumentProcessingError should be raised here or allow partial processing.

        full_text = "\n\n".join(text_content_parts)
        if options.get("clean_text", True):
            full_text = self._common_text_cleaning(full_text) or ""

        return {
            "text_content": full_text,
            "extracted_metadata": extracted_meta,
            "page_count": page_count,
            "processing_notes": notes,
            "image_ocr_results": image_ocr_res if image_ocr_res else None,
        }

    def _sync_process_docx(
        self, file_path: Path, options: Dict[str, Any]
    ) -> Dict[str, Any]:
        file_logger.info(f"Processing DOCX (sync): {file_path.name}")
        if not dep_manager.is_available("docx"):
            raise DocumentProcessingError(
                "python-docx unavailable.", file_path=file_path
            )

        text_parts: List[str] = []
        extracted_meta: Dict[str, Any] = {"source_format": "docx"}
        notes: List[str] = []
        tables_data: List[Dict[str, Any]] = []

        try:
            doc = DocxDocument(str(file_path))  # type: ignore
            cp = doc.core_properties
            for prop_name in [
                "title",
                "author",
                "subject",
                "keywords",
                "comments",
                "category",
                "last_modified_by",
            ]:
                val = getattr(cp, prop_name, None)
                if val:
                    extracted_meta[prop_name] = str(val)
            if cp.created:
                extracted_meta["created_date"] = cp.created.isoformat()
            if cp.modified:
                extracted_meta["modified_date"] = cp.modified.isoformat()
            if cp.revision:
                extracted_meta["revision"] = cp.revision

            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text.strip())

            if options.get("extract_tables", True) and doc.tables:
                notes.append(f"Extracting {len(doc.tables)} table(s).")
                for i, table in enumerate(doc.tables):
                    table_content = []
                    for row in table.rows:
                        row_content = [cell.text.strip() for cell in row.cells]
                        table_content.append(row_content)
                    tables_data.append({"table_index": i, "rows": table_content})
                    # Optionally add table text to main text_content (can make it messy)
                    # table_text_summary = "\n".join([" | ".join(row) for row in table_content[:5]]) # First 5 rows
                    # text_parts.append(f"\n--- TABLE {i+1} ---\n{table_text_summary}\n--- END TABLE ---\n")
        except DocxPackageNotFoundError:  # type: ignore
            notes.append("File is not a valid DOCX (Zip archive) package.")
            file_logger.error(f"Invalid DOCX package: {file_path.name}")
            raise DocumentProcessingError(
                f"Invalid DOCX file: {file_path.name}", file_path=file_path
            )
        except Exception as e:
            notes.append(f"Error processing DOCX: {str(e)}")
            file_logger.error(f"Error processing DOCX {file_path.name}.", exception=e)

        full_text = "\n\n".join(text_parts)
        if options.get("clean_text", True):
            full_text = self._common_text_cleaning(full_text) or ""

        return {
            "text_content": full_text,
            "extracted_metadata": extracted_meta,
            "processing_notes": notes,
            "tables_extracted": tables_data if tables_data else None,
            "table_count": len(tables_data),
        }

    def _sync_process_txt(
        self, file_path: Path, options: Dict[str, Any]
    ) -> Dict[str, Any]:
        file_logger.info(f"Processing TXT (sync): {file_path.name}")
        encodings_to_try = ["utf-8", "latin-1", "cp1252", "utf-16"]
        content = None
        encoding_used = "unknown"
        notes = []

        for enc in encodings_to_try:
            try:
                with open(file_path, "r", encoding=enc) as f:
                    content = f.read()
                encoding_used = enc
                notes.append(f"Read with encoding: {enc}")
                break
            except (UnicodeDecodeError, LookupError):
                continue  # LookupError for invalid encoding name

        if content is None:
            try:  # Fallback with error replacement
                with open(file_path, "rb") as fb:
                    raw_bytes = fb.read()
                content = raw_bytes.decode("utf-8", errors="replace")
                encoding_used = "utf-8-replace"
                notes.append("Read as UTF-8 replacing errors.")
            except Exception as e:
                raise DocumentProcessingError(
                    f"Could not read TXT file {file_path.name}: {str(e)}",
                    file_path=file_path,
                    cause=e,
                )

        if options.get("clean_text", True):
            content = self._common_text_cleaning(content) or ""
        return {
            "text_content": content,
            "extracted_metadata": {"encoding": encoding_used, "source_format": "txt"},
            "processing_notes": notes,
        }

    def _sync_process_markdown(
        self, file_path: Path, options: Dict[str, Any]
    ) -> Dict[str, Any]:
        file_logger.info(f"Processing Markdown (sync): {file_path.name}")
        notes = []
        text_content = ""
        if dep_manager.is_available("markdown") and dep_manager.is_available("bs4"):
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    md_content = f.read()
                html = markdown.markdown(md_content, extensions=["extra", "tables", "fenced_code"])  # type: ignore
                soup = BeautifulSoup(html, "html.parser")  # type: ignore
                # Extract text, trying to preserve some structure with newlines
                for element in soup.find_all(
                    ["p", "h1", "h2", "h3", "h4", "h5", "h6", "li", "pre", "code"]
                ):
                    text_content += element.get_text(separator=" ", strip=True) + "\n"
                if not text_content.strip():  # Fallback if above yields nothing
                    text_content = soup.get_text(separator="\n", strip=True)
                notes.append(
                    "Processed Markdown using python-markdown and BeautifulSoup."
                )
            except Exception as e:
                notes.append(
                    f"Error processing Markdown with libraries: {str(e)}. Falling back to text."
                )
                return self._sync_process_txt(
                    file_path, options
                )  # Fallback to plain text reading
        else:
            notes.append("Markdown/BeautifulSoup unavailable. Reading as plain text.")
            return self._sync_process_txt(file_path, options)

        if options.get("clean_text", True):
            text_content = self._common_text_cleaning(text_content) or ""
        return {
            "text_content": text_content,
            "extracted_metadata": {"source_format": "markdown"},
            "processing_notes": notes,
        }

    def _sync_process_html(
        self, file_path: Path, options: Dict[str, Any]
    ) -> Dict[str, Any]:
        file_logger.info(f"Processing HTML (sync): {file_path.name}")
        notes = []
        text_content = ""
        if not dep_manager.is_available("bs4"):
            raise DocumentProcessingError(
                "BeautifulSoup4 unavailable for HTML.", file_path=file_path
            )

        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                html_content = f.read()
            soup = BeautifulSoup(html_content, "html.parser")  # type: ignore

            # Remove script and style elements
            for script_or_style in soup(["script", "style"]):
                script_or_style.decompose()

            # Get text, trying to preserve paragraphs with newlines
            text_parts = []
            for element in soup.find_all(
                ["p", "h1", "h2", "h3", "h4", "h5", "h6", "div", "li", "td", "th"]
            ):
                block_text = element.get_text(separator=" ", strip=True)
                if block_text:
                    text_parts.append(block_text)
            text_content = "\n\n".join(text_parts)

            if not text_content.strip():  # Fallback
                text_content = soup.get_text(separator="\n", strip=True)

            extracted_meta = {
                "source_format": "html",
                "title": soup.title.string if soup.title else None,
            }
            notes.append("Processed HTML using BeautifulSoup.")
        except Exception as e:
            notes.append(f"Error processing HTML: {str(e)}. Reading as plain text.")
            file_logger.error(f"Error processing HTML {file_path.name}.", exception=e)
            return self._sync_process_txt(file_path, options)  # Fallback

        if options.get("clean_text", True):
            text_content = self._common_text_cleaning(text_content) or ""
        return {
            "text_content": text_content,
            "extracted_metadata": extracted_meta,
            "processing_notes": notes,
        }

    def _sync_process_rtf(
        self, file_path: Path, options: Dict[str, Any]
    ) -> Dict[str, Any]:
        file_logger.info(f"Processing RTF (sync): {file_path.name}")
        notes = []
        text_content = ""
        if dep_manager.is_available("striprtf"):
            try:
                with open(
                    file_path, "r", encoding="ascii", errors="ignore"
                ) as f:  # RTF typically ASCII/extended ASCII
                    rtf_content = f.read()
                text_content = rtf_to_text(rtf_content)  # type: ignore
                notes.append("Processed RTF using striprtf library.")
            except Exception as e:
                notes.append(
                    f"Error processing RTF with striprtf: {str(e)}. Fallback to text."
                )
                return self._sync_process_txt(file_path, options)
        else:
            notes.append(
                "striprtf library not available for RTF. Reading as plain text (may include RTF markup)."
            )
            return self._sync_process_txt(file_path, options)

        if options.get("clean_text", True):
            text_content = self._common_text_cleaning(text_content) or ""
        return {
            "text_content": text_content,
            "extracted_metadata": {"source_format": "rtf"},
            "processing_notes": notes,
        }

    def _sync_process_excel(
        self, file_path: Path, options: Dict[str, Any]
    ) -> Dict[str, Any]:
        file_logger.info(f"Processing Excel (sync): {file_path.name}")
        notes = []
        all_text_parts: List[str] = []
        tables_data: List[Dict[str, Any]] = []
        extracted_meta = {"source_format": "excel", "sheet_names": []}

        if not dep_manager.is_available("pandas"):
            raise DocumentProcessingError(
                "pandas unavailable for Excel.", file_path=file_path
            )
        # openpyxl for .xlsx, xlrd for .xls are soft dependencies of pandas for these formats

        try:
            xls = pd.ExcelFile(str(file_path))  # type: ignore
            extracted_meta["sheet_names"] = xls.sheet_names

            for sheet_idx, sheet_name in enumerate(xls.sheet_names):
                try:
                    df = xls.parse(sheet_name)
                    if df.empty:
                        notes.append(f"Sheet '{sheet_name}' is empty.")
                        continue

                    # Extract text content from sheet
                    sheet_text_parts = []
                    for col in df.columns:
                        # Convert column to string, handling potential mixed types robustly
                        try:
                            col_text = df[col].astype(str).str.strip().dropna()
                            sheet_text_parts.extend(col_text[col_text != ""].tolist())
                        except (
                            Exception
                        ):  # Catch errors during astype or string operations
                            notes.append(
                                f"Could not process column '{col}' in sheet '{sheet_name}' fully as text."
                            )

                    sheet_full_text = " ".join(sheet_text_parts)
                    if sheet_full_text.strip():
                        all_text_parts.append(
                            f"\n--- Sheet: {sheet_name} ---\n{sheet_full_text}"
                        )

                    # Extract table structure if requested
                    if options.get("extract_tables", True):
                        # Convert DataFrame to list of lists (header + rows)
                        table_as_list = [df.columns.tolist()] + df.values.astype(
                            str
                        ).tolist()
                        tables_data.append(
                            {
                                "sheet_name": sheet_name,
                                "table_index_in_doc": sheet_idx,
                                "rows": table_as_list,
                            }
                        )
                    notes.append(
                        f"Processed sheet: {sheet_name} ({len(df)} rows, {len(df.columns)} cols)."
                    )
                except Exception as sheet_e:
                    notes.append(
                        f"Error processing sheet '{sheet_name}': {str(sheet_e)}"
                    )
                    file_logger.warning(
                        f"Error processing sheet '{sheet_name}' in {file_path.name}.",
                        exception=sheet_e,
                    )

        except Exception as e:
            notes.append(f"Error processing Excel file: {str(e)}")
            file_logger.error(f"Error processing Excel {file_path.name}.", exception=e)
            # If it fails catastrophically, might not have any text.

        full_text = "\n\n".join(all_text_parts)
        if options.get("clean_text", True):
            full_text = self._common_text_cleaning(full_text) or ""

        return {
            "text_content": (
                full_text
                if options.get("processing_strategy_used")
                == ProcessingStrategy.FULL_PROCESSING
                else None
            ),  # Text only if full processing
            "extracted_metadata": extracted_meta,
            "processing_notes": notes,
            "tables_extracted": tables_data if tables_data else None,
            "table_count": len(tables_data),
        }

    def _sync_process_csv(
        self, file_path: Path, options: Dict[str, Any]
    ) -> Dict[str, Any]:
        file_logger.info(f"Processing CSV (sync): {file_path.name}")
        notes = []
        all_text_parts: List[str] = []
        tables_data: List[Dict[str, Any]] = []
        extracted_meta = {"source_format": "csv"}

        if not dep_manager.is_available("pandas"):
            raise DocumentProcessingError(
                "pandas unavailable for CSV.", file_path=file_path
            )

        try:
            # Try to infer encoding, common ones for CSV
            encodings_to_try = ["utf-8", "latin-1", "cp1252"]
            df = None
            for enc in encodings_to_try:
                try:
                    df = pd.read_csv(str(file_path), encoding=enc, low_memory=False)  # type: ignore
                    extracted_meta["encoding_detected"] = enc
                    notes.append(f"Read CSV with encoding: {enc}")
                    break
                except UnicodeDecodeError:
                    continue
                except pd.errors.ParserError as pe:  # type: ignore
                    notes.append(
                        f"Pandas parsing error with encoding {enc}: {str(pe)}. Trying next."
                    )

            if df is None:  # If all fail, try with error replacement
                try:
                    df = pd.read_csv(str(file_path), encoding="utf-8", errors="replace", low_memory=False)  # type: ignore
                    extracted_meta["encoding_detected"] = "utf-8-replace"
                    notes.append("Read CSV as UTF-8 replacing errors.")
                except Exception as final_e:
                    raise DocumentProcessingError(
                        f"Failed to parse CSV file {file_path.name}: {str(final_e)}",
                        file_path=file_path,
                        cause=final_e,
                    )

            if df.empty:
                notes.append("CSV file is empty or contains no data.")
            else:
                # Extract text content (concatenate all string cells)
                for col in df.columns:
                    try:
                        col_text = df[col].astype(str).str.strip().dropna()
                        all_text_parts.extend(col_text[col_text != ""].tolist())
                    except Exception:
                        notes.append(f"Could not process column '{col}' fully as text.")

                if options.get("extract_tables", True):
                    table_as_list = [df.columns.tolist()] + df.values.astype(
                        str
                    ).tolist()
                    tables_data.append(
                        {"table_index_in_doc": 0, "rows": table_as_list}
                    )  # CSV is one table
                notes.append(f"Processed CSV: {len(df)} rows, {len(df.columns)} cols.")

        except Exception as e:
            notes.append(f"Error processing CSV file: {str(e)}")
            file_logger.error(f"Error processing CSV {file_path.name}.", exception=e)

        full_text = " ".join(all_text_parts)
        if options.get("clean_text", True):
            full_text = self._common_text_cleaning(full_text) or ""

        return {
            "text_content": (
                full_text
                if options.get("processing_strategy_used")
                == ProcessingStrategy.FULL_PROCESSING
                else None
            ),
            "extracted_metadata": extracted_meta,
            "processing_notes": notes,
            "tables_extracted": tables_data if tables_data else None,
            "table_count": len(tables_data),
        }

    def _sync_process_powerpoint(
        self, file_path: Path, options: Dict[str, Any]
    ) -> Dict[str, Any]:
        file_logger.info(f"Processing PowerPoint (sync): {file_path.name}")
        notes = []
        text_parts: List[str] = []
        slide_count = 0
        extracted_meta = {"source_format": "pptx"}

        if not dep_manager.is_available("pptx"):
            raise DocumentProcessingError(
                "python-pptx unavailable.", file_path=file_path
            )

        try:
            prs = Presentation(str(file_path))  # type: ignore
            slide_count = len(prs.slides)
            # Extract core properties
            cp = prs.core_properties
            for prop_name in [
                "title",
                "author",
                "subject",
                "keywords",
                "comments",
                "category",
                "last_modified_by",
            ]:
                val = getattr(cp, prop_name, None)
                if val:
                    extracted_meta[prop_name] = str(val)
            if cp.created:
                extracted_meta["created_date"] = cp.created.isoformat()
            if cp.modified:
                extracted_meta["modified_date"] = cp.modified.isoformat()

            for i, slide in enumerate(prs.slides):
                slide_texts = [f"--- Slide {i+1} ---"]
                if slide.has_title and slide.shapes.title:
                    slide_texts.append(f"Title: {slide.shapes.title.text.strip()}")
                for shape in slide.shapes:
                    if (
                        hasattr(shape, "text_frame")
                        and shape.text_frame
                        and shape.text_frame.text.strip()
                    ):
                        # Avoid re-adding title if it's in a shape
                        if not (slide.has_title and shape == slide.shapes.title):
                            slide_texts.append(shape.text_frame.text.strip())
                    elif (
                        hasattr(shape, "has_table")
                        and shape.has_table
                        and options.get("extract_tables", True)
                    ):
                        # Basic table text extraction (no structure preserved here, add to notes)
                        table = shape.table
                        table_text_content = []
                        for row in table.rows:
                            row_cells = [cell.text.strip() for cell in row.cells]
                            table_text_content.append(" | ".join(row_cells))
                        if table_text_content:
                            slide_texts.append(
                                f"Table Content: {'; '.join(table_text_content)}"
                            )
                text_parts.append("\n".join(slide_texts))
            notes.append(f"Processed {slide_count} slides.")
        except PptxPackageNotFoundError:  # type: ignore
            notes.append(
                f"File {file_path.name} is not a valid PowerPoint (Zip) package."
            )
            file_logger.error(f"Invalid PowerPoint package: {file_path.name}")
            raise DocumentProcessingError(
                f"Invalid PowerPoint file: {file_path.name}", file_path=file_path
            )
        except Exception as e:
            notes.append(f"Error processing PowerPoint: {str(e)}")
            file_logger.error(
                f"Error processing PowerPoint {file_path.name}.", exception=e
            )

        full_text = "\n\n".join(text_parts)
        if options.get("clean_text", True):
            full_text = self._common_text_cleaning(full_text) or ""

        return {
            "text_content": full_text,
            "extracted_metadata": extracted_meta,
            "page_count": slide_count,  # Use page_count for slides
            "processing_notes": notes,
        }

    def _sync_process_image(
        self, file_path: Path, options: Dict[str, Any]
    ) -> Dict[str, Any]:
        file_logger.info(f"Processing Image for OCR (sync): {file_path.name}")
        notes = []
        text_content = ""
        extracted_meta: Dict[str, Any] = {"source_format": "image"}
        file_hash = ocr_cache.compute_file_hash(file_path)
        cached = ocr_cache.get(file_path, page=1, file_hash=file_hash)

        if not (
            dep_manager.is_available("pytesseract") and dep_manager.is_available("PIL")
        ):
            raise DocumentProcessingError(
                "Pytesseract/Pillow unavailable for image OCR.", file_path=file_path
            )

        try:
            pil_img = Image.open(str(file_path))  # type: ignore
            meta_info: Dict[str, Any] = {
                "format": pil_img.format,
                "mode": pil_img.mode,
                "size": pil_img.size,
                "info": pil_img.info,
            }
            extracted_meta.update(meta_info)

            if cached and not options.get("force_ocr", False):
                text_content = cached
                notes.append("OCR text loaded from cache.")
            else:
                text_content = pytesseract.image_to_string(
                    pil_img,
                    lang=options.get("ocr_language", "eng"),
                    config=f"--dpi {options.get('ocr_dpi',300)}",
                )  # type: ignore
                text_content = text_content.strip()
                if text_content:
                    notes.append(
                        f"OCR successful, extracted {len(text_content)} characters."
                    )
                    ocr_cache.set(file_path, 1, text_content, file_hash=file_hash)
                else:
                    notes.append("OCR performed, but no text detected.")
        except UnidentifiedImageError:  # type: ignore
            notes.append(f"Cannot identify image file: {file_path.name}")
            file_logger.error(f"Pillow UnidentifiedImageError for {file_path.name}")
            raise DocumentProcessingError(
                f"Cannot identify image file: {file_path.name}", file_path=file_path
            )
        except pytesseract.TesseractNotFoundError:  # type: ignore
            notes.append("Tesseract OCR engine not found or not in PATH.")
            file_logger.error(
                "TesseractNotFoundError. Ensure Tesseract is installed and in PATH."
            )
            raise DocumentProcessingError(
                "Tesseract OCR not installed/configured correctly.", file_path=file_path
            )
        except Exception as e:
            notes.append(f"Error during image OCR: {str(e)}")
            file_logger.error(
                f"Error during image OCR for {file_path.name}.", exception=e
            )
            # Continue, text_content will be empty.

        if options.get("clean_text", True) and text_content:
            text_content = self._common_text_cleaning(text_content) or ""

        return {
            "text_content": text_content,
            "extracted_metadata": extracted_meta,
            "processing_notes": notes,
        }

    # ---- Additional Formats ----

    def _process_eml_async(self, file_path: Path, options: Dict[str, Any]) -> Dict[str, Any]:
        return self._sync_process_eml(file_path, options)

    def _sync_process_eml(self, file_path: Path, options: Dict[str, Any]) -> Dict[str, Any]:
        file_logger.info(f"Processing EML (sync): {file_path.name}")
        notes: List[str] = []
        text_content = ""
        metadata: Dict[str, Any] = {"source_format": "eml"}
        try:
            with open(file_path, "rb") as f:
                msg = email.message_from_binary_file(f)
            metadata.update({
                "from": msg.get("From"),
                "to": msg.get("To"),
                "subject": msg.get("Subject"),
                "date": msg.get("Date"),
            })
            attachments: List[str] = []
            body_parts: List[str] = []
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_maintype() == "multipart":
                        continue
                    disp = part.get("Content-Disposition", "")
                    if disp.startswith("attachment"):
                        if part.get_filename():
                            attachments.append(part.get_filename())
                        continue
                    if part.get_content_type() == "text/plain":
                        charset = part.get_content_charset() or "utf-8"
                        payload = part.get_payload(decode=True)
                        if payload:
                            body_parts.append(payload.decode(charset, errors="ignore"))
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    charset = msg.get_content_charset() or "utf-8"
                    body_parts.append(payload.decode(charset, errors="ignore"))
            if attachments:
                metadata["attachments"] = attachments
            text_content = "\n".join(body_parts)
        except Exception as e:
            notes.append(f"Error parsing email: {str(e)}")
            file_logger.error(f"Error processing EML {file_path.name}.", exception=e)

        if options.get("clean_text", True) and text_content:
            text_content = self._common_text_cleaning(text_content) or ""

        return {
            "text_content": text_content,
            "extracted_metadata": metadata,
            "processing_notes": notes,
        }

    def _process_zip_async(self, file_path: Path, options: Dict[str, Any]) -> Dict[str, Any]:
        return self._sync_process_zip(file_path, options)

    def _sync_process_zip(self, file_path: Path, options: Dict[str, Any]) -> Dict[str, Any]:
        file_logger.info(f"Processing ZIP (sync): {file_path.name}")
        notes: List[str] = []
        extracted_meta: Dict[str, Any] = {"source_format": "zip"}
        files_info: List[Dict[str, Any]] = []
        try:
            with zipfile.ZipFile(file_path, "r") as zf:
                for info in zf.infolist():
                    files_info.append({"name": info.filename, "size": info.file_size})
            extracted_meta["contained_files"] = files_info
            notes.append(f"Zip contains {len(files_info)} entries.")
        except zipfile.BadZipFile:
            notes.append("Invalid ZIP archive.")
            file_logger.error(f"BadZipFile: {file_path.name}")
            raise DocumentProcessingError("Invalid ZIP archive", file_path=file_path)
        except Exception as e:
            notes.append(f"Error processing ZIP: {str(e)}")
            file_logger.error(f"Error processing ZIP {file_path.name}.", exception=e)

        return {"extracted_metadata": extracted_meta, "processing_notes": notes}

    def _process_json_async(self, file_path: Path, options: Dict[str, Any]) -> Dict[str, Any]:
        return self._sync_process_json(file_path, options)

    def _sync_process_json(self, file_path: Path, options: Dict[str, Any]) -> Dict[str, Any]:
        file_logger.info(f"Processing JSON (sync): {file_path.name}")
        notes: List[str] = []
        text_content = ""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            text_content = json.dumps(data, indent=2)
        except Exception as e:
            notes.append(f"Error parsing JSON: {str(e)}")
            file_logger.error(f"Error processing JSON {file_path.name}.", exception=e)

        if options.get("clean_text", True) and text_content:
            text_content = self._common_text_cleaning(text_content) or ""

        return {
            "text_content": text_content,
            "extracted_metadata": {"source_format": "json"},
            "processing_notes": notes,
        }

    def _process_yaml_async(self, file_path: Path, options: Dict[str, Any]) -> Dict[str, Any]:
        return self._sync_process_yaml(file_path, options)

    def _sync_process_yaml(self, file_path: Path, options: Dict[str, Any]) -> Dict[str, Any]:
        file_logger.info(f"Processing YAML (sync): {file_path.name}")
        notes: List[str] = []
        text_content = ""
        if not dep_manager.is_available("yaml"):
            notes.append("PyYAML unavailable.")
            raise DocumentProcessingError("PyYAML not installed", file_path=file_path)
        try:
            import yaml  # type: ignore
            with open(file_path, "r", encoding="utf-8") as f:
                data = yaml.safe_load(f)
            text_content = yaml.safe_dump(data)
        except Exception as e:
            notes.append(f"Error parsing YAML: {str(e)}")
            file_logger.error(f"Error processing YAML {file_path.name}.", exception=e)

        if options.get("clean_text", True) and text_content:
            text_content = self._common_text_cleaning(text_content) or ""

        return {
            "text_content": text_content,
            "extracted_metadata": {"source_format": "yaml"},
            "processing_notes": notes,
        }

    async def health_check(self) -> Dict[str, Any]:  # Public method
        base_status = await super().health_check()
        base_status["agent_name"] = self.name
        base_status["config_summary"] = self.get_config_summary_params()

        supported_formats_with_deps: List[str] = []
        unsupported_due_to_deps: List[str] = []
        for content_type, config_details in self.file_type_configs.items():
            deps_ok = all(
                dep_manager.is_available(dep) for dep in config_details.get("deps", [])
            )
            if deps_ok:
                supported_formats_with_deps.append(content_type.name)
            else:
                unsupported_due_to_deps.append(content_type.name)

        base_status["supported_formats"] = supported_formats_with_deps
        if unsupported_due_to_deps:
            base_status["unsupported_formats_missing_deps"] = unsupported_due_to_deps

        return base_status
